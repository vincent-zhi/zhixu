"""Assemble creative-mode PNGs into a single 16:9 PPTX, one image per slide, full-bleed.

Usage:
    python3 build_pptx.py --deck-dir <deck_dir> [--output <path>]

Default output: <deck_dir>/<deck_id>.pptx (deck_id from task_pack.json).

Behavior:
- Reads task_pack.json for page_count and deck_id.
- Reads outline.json (if present) to order pages by page_no; otherwise falls back
  to lexical order of pages/page_*.png.
- Missing PNGs are skipped with a warning line on stderr; they do NOT abort.
- Each slide is 16:9 (13.333 x 7.5 inches). Images are stretched to fill the slide.
- Emits a single JSON line on stdout summarizing the result.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def resolve_page_order(deck: Path, page_count: int) -> list[tuple[int, Path]]:
    outline_path = deck / "outline.json"
    pages_dir = deck / "pages"
    ordered: list[tuple[int, Path]] = []
    if outline_path.exists():
        data = json.loads(outline_path.read_text(encoding="utf-8"))
        for page in data.get("pages", []):
            n = int(page["page_no"])
            png = pages_dir / f"page_{n:03d}.png"
            ordered.append((n, png))
        ordered.sort(key=lambda x: x[0])
    else:
        for n in range(1, page_count + 1):
            ordered.append((n, pages_dir / f"page_{n:03d}.png"))
    return ordered


def build(deck: Path, output: Path) -> dict:
    tp = json.loads((deck / "task_pack.json").read_text(encoding="utf-8"))
    deck_id = tp.get("deck_id") or deck.name
    page_count = int(tp["params"]["page_count"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    included: list[int] = []
    missing: list[int] = []
    for n, png in resolve_page_order(deck, page_count):
        slide = prs.slides.add_slide(blank_layout)
        if png.exists():
            slide.shapes.add_picture(str(png), 0, 0, width=SLIDE_W, height=SLIDE_H)
            included.append(n)
        else:
            missing.append(n)
            print(f"[build_pptx] page_{n:03d}.png missing, inserting blank slide", file=sys.stderr)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    return {
        "deck_id": deck_id,
        "output": str(output),
        "total_slides": len(included) + len(missing),
        "included_pages": included,
        "missing_pages": missing,
    }


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser()
    p.add_argument("--deck-dir", type=Path, required=True)
    p.add_argument("--output", type=Path, default=None,
                   help="Output .pptx path. Default: <deck_dir>/<deck_id>.pptx")
    args = p.parse_args(argv)

    deck = args.deck_dir.expanduser().resolve()
    if not (deck / "task_pack.json").exists():
        print(json.dumps({"error": f"task_pack.json missing in {deck}"}, ensure_ascii=False))
        return 2

    if args.output is None:
        tp = json.loads((deck / "task_pack.json").read_text(encoding="utf-8"))
        deck_id = tp.get("deck_id") or deck.name
        output = deck / f"{deck_id}.pptx"
    else:
        output = args.output.expanduser().resolve()

    result = build(deck, output)
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
